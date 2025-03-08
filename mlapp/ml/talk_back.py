#!/usr/bin/env python
import os
import sys
import json
from pptx import Presentation
import google.generativeai as genai

def extract_text_from_ppt(ppt_path, output_file=None):
    """
    Extract text from a PowerPoint file.
    
    Args:
        ppt_path (str): Path to the PowerPoint file.
        output_file (str, optional): Path to save the extracted text. If None, prints to console.
    
    Returns:
        str: Extracted text.
    """
    if not os.path.exists(ppt_path):
        raise FileNotFoundError(f"The file {ppt_path} does not exist.")
    
    # Open presentation
    prs = Presentation(ppt_path)
    
    # Initialize text variable
    all_text = []
    
    # Iterate through slides
    for i, slide in enumerate(prs.slides):
        slide_text = [f"\n\n--- Slide {i+1} ---\n"]
        
        # Extract text from shapes
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_text.append(shape.text)
                
        # Extract text from tables
        for shape in slide.shapes:
            if hasattr(shape, "has_table") and shape.has_table:
                for row in shape.table.rows:
                    row_text = [cell.text for cell in row.cells if cell.text]
                    if row_text:
                        slide_text.append(" | ".join(row_text))
        
        all_text.extend(slide_text)
    
    # Join all text
    full_text = "\n".join(all_text)
    
    # Output handling
    if output_file:
        with open(output_file, 'w', encoding='utf-8') as f:
            f.write(full_text)
        print(f"Text extracted and saved to {output_file}")
    
    return full_text

def organize_text_with_llm(text, topic="gravity", api_key="AIzaSyCvIuLul5DP1QmaFh1QFVN0_-acLlKP1LI", model_name="gemini-2.0-flash"):
    # Configure the API
    genai.configure(api_key=api_key)
    
    # Create model instance
    model = genai.GenerativeModel(model_name)
    
    # Create prompt
    prompt = f"""
    Analyze this text about {topic} and extract:
    1. Key concepts (7-10 important terms)
    2. Required concepts (3-4 essential terms)
    3. 3-5 hint questions to guide understanding
    4. A concise good explanation (1-3 sentences)
    
    Return ONLY a valid JSON with this structure:
    {{
        "science": {{
            "{topic}": {{
                "key_concepts": ["concept1", "concept2", ...],
                "required_concepts": ["required1", "required2", ...],
                "hint_questions": ["question1?", "question2?", ...],
                "example_good_explanation": "Concise explanation here."
            }}
        }}
    }}
    """
    
    # Get response from model
    response = model.generate_content(prompt + "\n\nText to analyze:\n" + text)
    
    # Extract JSON from response
    try:
        # Try to parse the response directly as JSON
        result = json.loads(response.text)
        return result
    except json.JSONDecodeError:
        # If not direct JSON, search for JSON block in text
        import re
        json_pattern = r'```(?:json)?\s*([\s\S]*?)\s*```'
        match = re.search(json_pattern, response.text)
        if match:
            try:
                return json.loads(match.group(1))
            except:
                # If still fails, try to extract everything between curly braces
                full_json = re.search(r'({[\s\S]*})', response.text)
                if full_json:
                    return json.loads(full_json.group(1))
        
        # If all extraction methods fail
        raise ValueError("Could not extract valid JSON from model response")

def main():
    if len(sys.argv) < 4:
        print("Usage: python ppt_analyzer.py <ppt_file_path> <topic> <output_json_path>")
        sys.exit(1)
    
    ppt_path = sys.argv[1]
    topic = sys.argv[2]
    output_json_path = sys.argv[3]
    
    # Extract text from PowerPoint
    text_output_path = os.path.splitext(ppt_path)[0] + ".txt"
    extracted_text = extract_text_from_ppt(ppt_path, text_output_path)
    
    # Analyze the extracted text with LLM
    result = organize_text_with_llm(extracted_text, topic)
    
    # Save the result to JSON
    with open(output_json_path, "w") as f:
        json.dump(result, f, indent=4)
    
    print(f"Successfully saved PowerPoint analysis to {output_json_path}")
    print(json.dumps(result, indent=2))

if __name__ == "__main__":
    main()

# Example of how to integrate with Django:
"""
# In your Django views.py:

from django.shortcuts import render
from django.http import HttpResponse, JsonResponse
from django.conf import settings
import os
from .ppt_analyzer import extract_text_from_ppt, organize_text_with_llm

def analyze_ppt_view(request):
    if request.method == 'POST' and request.FILES.get('ppt_file'):
        # Get the uploaded file and topic
        ppt_file = request.FILES['ppt_file']
        topic = request.POST.get('topic', 'general')
        
        # Save the uploaded file temporarily
        temp_path = os.path.join(settings.MEDIA_ROOT, 'temp', ppt_file.name)
        os.makedirs(os.path.dirname(temp_path), exist_ok=True)
        
        with open(temp_path, 'wb+') as destination:
            for chunk in ppt_file.chunks():
                destination.write(chunk)
        
        # Extract text
        text_output_path = os.path.join(settings.MEDIA_ROOT, 'extracted', 
                                        os.path.splitext(ppt_file.name)[0] + '.txt')
        os.makedirs(os.path.dirname(text_output_path), exist_ok=True)
        
        extracted_text = extract_text_from_ppt(temp_path, text_output_path)
        
        # Analyze with LLM
        result = organize_text_with_llm(extracted_text, topic)
        
        # Save result to JSON
        json_output_path = os.path.join(settings.MEDIA_ROOT, 'analyzed', 
                                        os.path.splitext(ppt_file.name)[0] + '.json')
        os.makedirs(os.path.dirname(json_output_path), exist_ok=True)
        
        with open(json_output_path, 'w') as f:
            json.dump(result, f, indent=4)
        
        # Clean up temporary file
        os.remove(temp_path)
        
        # Return the result as JSON response
        return JsonResponse(result)
    
    return render(request, 'upload_ppt_for_analysis.html')
"""