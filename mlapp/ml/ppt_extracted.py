#!/usr/bin/env python
import os
import sys
from pptx import Presentation

def extract_text_from_ppt(ppt_path, output_file=None):
    """
    Extract text from a PowerPoint file.
    
    Args:
        ppt_path (str): Path to the PowerPoint file.
        output_file (str, optional): Path to save the extracted text. If None, prints to console.
    
    Returns:
        str: Extracted text.
    """
    if not os.path.exists(ppt_path):
        raise FileNotFoundError(f"The file {ppt_path} does not exist.")
    
    # Open presentation
    prs = Presentation(ppt_path)
    
    # Initialize text variable
    all_text = []
    
    # Iterate through slides
    for i, slide in enumerate(prs.slides):
        slide_text = [f"\n\n--- Slide {i+1} ---\n"]
        
        # Extract text from shapes
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_text.append(shape.text)
                
        # Extract text from tables
        for shape in slide.shapes:
            if hasattr(shape, "has_table") and shape.has_table:
                for row in shape.table.rows:
                    row_text = [cell.text for cell in row.cells if cell.text]
                    if row_text:
                        slide_text.append(" | ".join(row_text))
        
        all_text.extend(slide_text)
    
    # Join all text
    full_text = "\n".join(all_text)
    
    # Output handling
    if output_file:
        with open(output_file, 'w', encoding='utf-8') as f:
            f.write(full_text)
        print(f"Text extracted and saved to {output_file}")
    else:
        print(full_text)
    
    return full_text

def main():
    # Check if correct number of arguments
    if len(sys.argv) < 2:
        print("Usage: python ppt_extractor.py <ppt_file_path> [output_file_path]")
        sys.exit(1)
    
    # Get file paths from command line arguments
    ppt_path = sys.argv[1]
    output_path = sys.argv[2] if len(sys.argv) > 2 else None
    
    # Extract text
    extract_text_from_ppt(ppt_path, output_path)

if __name__ == "__main__":
    main()

# Example of how to integrate with Django:
"""
# In your Django view:

from django.shortcuts import render
from django.http import HttpResponse, FileResponse
from django.conf import settings
import os
from .ppt_extractor import extract_text_from_ppt

def extract_ppt_view(request):
    if request.method == 'POST' and request.FILES.get('ppt_file'):
        # Get the uploaded file
        ppt_file = request.FILES['ppt_file']
        
        # Save the uploaded file temporarily
        temp_path = os.path.join(settings.MEDIA_ROOT, 'temp', ppt_file.name)
        os.makedirs(os.path.dirname(temp_path), exist_ok=True)
        
        with open(temp_path, 'wb+') as destination:
            for chunk in ppt_file.chunks():
                destination.write(chunk)
        
        # Define output path
        output_filename = os.path.splitext(ppt_file.name)[0] + '.txt'
        output_path = os.path.join(settings.MEDIA_ROOT, 'extracted', output_filename)
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        
        # Extract text
        extract_text_from_ppt(temp_path, output_path)
        
        # Clean up temporary file
        os.remove(temp_path)
        
        # Return the extracted text file as a download
        response = FileResponse(open(output_path, 'rb'))
        response['Content-Disposition'] = f'attachment; filename="{output_filename}"'
        return response
    
    return render(request, 'upload_ppt.html')
"""